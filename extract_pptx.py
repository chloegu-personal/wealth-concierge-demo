import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text(pptx_path, output_path):
    try:
        def get_text_from_shape(shape):
            text = ""
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            if shape.shape_type == 6:  # Group shape
                for s in shape.shapes:
                    text += get_text_from_shape(s)
            return text

        prs = Presentation(pptx_path)
        with open(output_path, "w", encoding="utf-8") as f:
            for i, slide in enumerate(prs.slides):
                f.write(f"--- Slide {i+1} ---\n")
                for shape in slide.shapes:
                    f.write(get_text_from_shape(shape))
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    f.write(f"--- Notes ---\n{notes}\n")
                f.write("\n")
                
                # Extract images from specific slides
                if i + 1 in [24, 25]:
                    if not os.path.exists("extracted_images"):
                        os.makedirs("extracted_images")
                    for shape_idx, shape in enumerate(slide.shapes):
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image = shape.image
                            image_bytes = image.blob
                            image_filename = f"extracted_images/slide_{i+1}_image_{shape_idx}.{image.ext}"
                            with open(image_filename, "wb") as img_file:
                                img_file.write(image_bytes)
                            print(f"Extracted image: {image_filename}")

        print(f"Success: Extracted text to {output_path}")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    extract_text("SC4025 Digital Product Management - 13Weeks - Rapid MVP-1.pptx", "extracted_text_utf8.txt")
